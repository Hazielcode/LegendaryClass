# -*- coding: utf-8 -*-
"""
Genera PDF, DOCX y PPTX del Plan de Marketing - LegendaryClass
Estilo institucional TECSUP: azul oscuro #1F4D78 / azul medio #2E74B5
"""

import os, io
BASE = "C:/xampp/htdocs/LegendaryClass"
LOGO = os.path.join(BASE, "portada_ref.png")

# ─────────────────────────────────────────────────────────────────
# CONTENIDO COMPARTIDO
# ─────────────────────────────────────────────────────────────────
INTEGRANTES = [
    "Aguirre Saavedra, Juan Alexis",
    "Alfonso Solorzano, Samir Haziel",
    "Galvan Morales, Luis Enrique",
    "Galvan Guerrero, Matias",
]
CURSO   = "Marketing y Comercialización de Nuevos Productos"
DOCENTE = "Mg. Miluska Horna Elera"
CARRERA = "Diseño y Desarrollo de Software"
ANIO    = "Lima – Perú 2026"

CAP1 = {
    "titulo": "CAPÍTULO 1: INTRODUCCIÓN",
    "secciones": [
        ("Nombre de la Empresa", "LegendaryClass"),
        ("Misión",
         "Transformar la experiencia educativa en instituciones escolares mediante una plataforma "
         "de gamificación que motive a los estudiantes, empodere a los docentes y mantenga a las "
         "familias conectadas con el progreso académico, convirtiendo el aprendizaje cotidiano en "
         "una aventura épica."),
        ("Visión",
         "Ser la plataforma de gamificación educativa líder en Latinoamérica para el 2030, "
         "presente en más de 500 instituciones educativas y reconocida por elevar el compromiso "
         "estudiantil y el rendimiento académico a través de la innovación tecnológica."),
        ("Valores",
         "• Innovación: Buscamos constantemente nuevas formas de hacer el aprendizaje más atractivo.\n"
         "• Compromiso: Nos dedicamos al éxito educativo de cada institución que confía en nuestra plataforma.\n"
         "• Accesibilidad: Soluciones adaptables a la realidad económica de los colegios peruanos.\n"
         "• Transparencia: Información clara y honesta sobre el progreso de cada estudiante.\n"
         "• Colaboración: Fomentamos el trabajo conjunto entre directores, maestros, alumnos y familias."),
        ("Objetivo General",
         "Diseñar una estrategia de marketing integral para el lanzamiento de LegendaryClass como "
         "plataforma de gamificación educativa B2B dirigida a instituciones de educación básica regular "
         "en el Perú, estableciendo modelos de comercialización sostenibles y escalables."),
        ("Objetivos Específicos",
         "1. Identificar y segmentar el mercado objetivo de colegios privados en Lima Metropolitana "
         "para establecer propuestas de valor diferenciadas según el tamaño y capacidad económica de cada institución.\n"
         "2. Definir una estrategia de precios dual (suscripción SaaS e implementación personalizada) "
         "que contemple los costos operativos reales y sea competitiva frente a alternativas extranjeras.\n"
         "3. Diseñar un plan de comunicación y promoción digital orientado a directores y coordinadores "
         "académicos que genere leads calificados durante el primer año de operación."),
        ("Descripción del Producto",
         "LegendaryClass es una plataforma web de gamificación educativa que transforma el aula "
         "tradicional en una experiencia de juego de rol (RPG). Los estudiantes crean personajes "
         "con clases únicas (Mago, Guerrero, Ninja, Arquero, Lanzador), acumulan puntos de "
         "experiencia, suben de nivel, completan misiones y canjean recompensas. Los docentes "
         "registran comportamientos positivos y negativos, gestionan aulas virtuales y generan "
         "reportes de rendimiento. Los directores monitorean toda la institución en tiempo real, "
         "y los padres pueden seguir el progreso académico y conductual de sus hijos. La plataforma "
         "es 100% web, responsive y en español, adaptable a la identidad visual de cada colegio."),
    ]
}

CAP2 = {
    "titulo": "CAPÍTULO 2: ANÁLISIS DEL MACROENTORNO",
    "subsecciones": [
        ("Factores Políticos",
         "• El MINEDU promueve la transformación digital mediante programas como 'Aprendo en Casa'.\n"
         "• Iniciativas de dotación tecnológica en colegios públicos (PRONIED).\n"
         "• La Ley N° 29733 regula el tratamiento de datos de menores de edad.\n"
         "• La inestabilidad política puede afectar decisiones de inversión en colegios públicos."),
        ("Factores Económicos",
         "• El mercado global de EdTech alcanzó USD 142 mil millones en 2023 (HolonIQ, 2023).\n"
         "• La recuperación post-pandemia incrementó la disposición de colegios privados a invertir en tecnología.\n"
         "• La fluctuación del tipo de cambio encarece alternativas extranjeras en USD.\n"
         "• NSE B y C: segmento de mayor volumen, sensible al precio; requiere modelos flexibles."),
        ("Factores Sociales",
         "• Generación Z y Alpha: nativos digitales habituados a videojuegos; alta receptividad a la gamificación.\n"
         "• Creciente desmotivación estudiantil en nivel secundaria a nivel nacional.\n"
         "• Padres urbanos demandan mayor visibilidad del comportamiento y rendimiento de sus hijos.\n"
         "• Brecha digital entre colegios urbanos y rurales limita la expansión inicial."),
        ("Factores Tecnológicos",
         "• Penetración de internet en Lima Metropolitana supera el 80% en hogares (INEI, 2023).\n"
         "• Expansión del uso de smartphones y tablets facilita el acceso a plataformas web responsive.\n"
         "• ClassDojo y Classcraft validan globalmente la demanda por herramientas de gamificación.\n"
         "• Stack moderno: Laravel 11, MongoDB, TailwindCSS, AlpineJS."),
        ("Factores Ecológicos",
         "• Solución 100% digital: elimina papel, tarjetas físicas y materiales de sistemas de recompensa tradicionales.\n"
         "• Infraestructura cloud permite optar por proveedores con certificaciones de sostenibilidad.\n"
         "• Posicionamiento eco-friendly como argumento adicional para instituciones con valores ambientales."),
        ("Factores Legales",
         "• Ley N° 29733: consentimiento obligatorio de padres/tutores para datos de menores de edad.\n"
         "• Formalización de empresa (SAC) con contratos SaaS que incluyan SLA y cláusulas de privacidad.\n"
         "• Registro de propiedad intelectual del software en INDECOPI.\n"
         "• Cumplimiento de Ley de Firmas y Certificados Digitales para contratos remotos."),
    ],
    "impacto": (
        "El contexto actual es altamente favorable para el lanzamiento de LegendaryClass. "
        "La digitalización educativa acelerada, la generación nativa digital como usuario final, "
        "y la debilidad de las alternativas extranjeras frente a una solución local en soles y en "
        "español crean una ventana de oportunidad clara. El principal riesgo a gestionar es el "
        "cumplimiento de la normativa de protección de datos de menores, que debe incorporarse al "
        "diseño del sistema antes del lanzamiento comercial."
    )
}

CAP3 = {
    "titulo": "CAPÍTULO 3: SEGMENTACIÓN Y PERFIL DEL CONSUMIDOR",
    "geo": [
        ["Primario",   "Lima Metropolitana", "Mayor concentración de colegios privados con infraestructura digital y capacidad de inversión"],
        ["Secundario", "Arequipa, Trujillo, Piura", "Ciudades con economía dinámica y colegios privados en crecimiento"],
        ["Terciario",  "Resto del país", "Expansión a mediano plazo, condicionada a conectividad"],
    ],
    "demo": (
        "Decisor de compra: directores/subdirectores 35–60 años, colegios privados 150–2,000 alumnos, NSE B/C.\n"
        "Usuario final: estudiantes 8–17 años, nativos digitales, habituados a mecánicas de videojuegos.\n"
        "Usuario operativo: docentes 25–50 años con acceso a dispositivo en el aula."
    ),
    "psico": (
        "• Instituciones que buscan diferenciarse frente a competidores escolares.\n"
        "• Directivos con mentalidad innovadora, apertura a nuevas metodologías.\n"
        "• Maestros frustrados con sistemas manuales (listas en papel, stickers).\n"
        "• Padres que valoran el seguimiento activo del desarrollo de sus hijos."
    ),
    "conductual": [
        ["Ocasión de uso",       "Diario (comportamientos), semanal (reportes), mensual (análisis institucional)"],
        ["Beneficios buscados",  "Motivación estudiantil, ahorro de tiempo, diferenciación institucional"],
        ["Estado de uso",        "Colegios que ya usan ClassDojo o Google Classroom (early adopters)"],
        ["Lealtad",              "Alta una vez adoptado, por datos acumulados y habituación del equipo"],
        ["Actitud",              "Colegios con historial de adopción de herramientas pedagógicas digitales"],
    ],
    "persona": (
        '"Director Innovador": Rodrigo, 45 años, director de colegio privado (400 alumnos, Surco). '
        'Busca mejorar disciplina sin aumentar carga docente. Ya usa Google Workspace pero necesita '
        'algo que enganche a los estudiantes. Decide con base en resultados concretos y referencias. '
        'Dispuesto a invertir S/. 400–700/mes si ve retorno en participación estudiantil.'
    )
}

CAP4 = {
    "titulo": "CAPÍTULO 4: ESTRATEGIAS DE MARKETING MIX",
    "producto_feats": [
        ["Sistema RPG completo (XP, niveles, personajes)", "Aumenta motivación y asistencia estudiantil"],
        ["Módulo de comportamientos positivos/negativos",  "Reemplaza sistemas manuales; ahorra tiempo al docente"],
        ["Panel del director con métricas en tiempo real", "Decisiones basadas en datos concretos"],
        ["Módulo para padres",                             "Mayor involucramiento familiar"],
        ["Importación masiva de alumnos (Excel)",          "Implementación rápida sin carga manual"],
        ["Tienda de recompensas personalizable",           "El colegio define sus propias recompensas"],
        ["100% en español, adaptado a Perú",               "Sin barreras de idioma ni de moneda"],
        ["Whitelabel: personalizable con marca del colegio", "El colegio lo siente propio, no genérico"],
    ],
    "comparativa": [
        ["Función",               "LegendaryClass", "ClassDojo", "Classcraft"],
        ["Sistema RPG / Niveles", "Sí – Completo",  "No tiene",  "Básico"],
        ["100% en español",       "Sí",             "Parcial",   "No – Inglés"],
        ["Precio en soles",       "Sí",             "Gratuito",  "No – USD"],
        ["Panel de Directores",   "Sí",             "No",        "No"],
        ["Módulo para Padres",    "Sí",             "Sí",        "No"],
        ["Tienda de Recompensas", "Sí – Completa",  "No",        "Limitado"],
    ],
    "propuesta": (
        '"LegendaryClass convierte el salón de clases en una experiencia épica de aprendizaje, '
        'incrementando el compromiso y la motivación estudiantil, mientras otorga a los docentes '
        'una herramienta ágil para gestionar comportamientos y reconocer el esfuerzo, '
        'todo en español y al precio de la realidad peruana."'
    ),
    "precio_saas": [
        ["Plan",         "N° de alumnos", "Precio mensual", "Precio anual (10% dcto.)"],
        ["Básico",       "Hasta 150",     "S/. 180",        "S/. 1,944"],
        ["Estándar",     "151 – 400",     "S/. 350",        "S/. 3,780"],
        ["Avanzado",     "401 – 800",     "S/. 580",        "S/. 6,264"],
        ["Institucional","+800",          "S/. 900",        "S/. 9,720"],
    ],
    "precio_vd": [
        ["Tipo de colegio",              "Precio único",  "Incluye"],
        ["Pequeño (hasta 200 alumnos)",  "S/. 3,500",    "Personalización básica + 3 meses soporte"],
        ["Mediano (201–600 alumnos)",    "S/. 7,500",    "Whitelabel + capacitaciones + 6 meses soporte"],
        ["Grande (+600 alumnos)",        "S/. 15,000",   "Desarrollo a medida + integración + 12 meses soporte"],
    ],
    "costos": [
        ["Concepto",                          "Costo mensual", "Costo anual"],
        ["Servidor VPS (hosting)",            "S/. 200",       "S/. 2,400"],
        ["Base de datos MongoDB Atlas",       "S/. 120",       "S/. 1,440"],
        ["Dominio + SSL",                     "S/. 5",         "S/. 60"],
        ["Herramientas de soporte",           "S/. 80",        "S/. 960"],
        ["Marketing digital (pauta)",         "S/. 500",       "S/. 6,000"],
        ["Materiales y eventos BTL",          "S/. 200",       "S/. 2,400"],
        ["Capacitaciones a clientes",         "S/. 150",       "S/. 1,800"],
        ["Contingencias y otros",             "S/. 100",       "S/. 1,200"],
        ["TOTAL",                             "S/. 1,355",     "S/. 16,260"],
    ],
}

# ═════════════════════════════════════════════════════════════════
# 1. DOCX  –  Formato APA 7ª edición
# ═════════════════════════════════════════════════════════════════
from docx import Document
from docx.shared import Pt, Cm, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import copy

AZUL_OSCURO = RGBColor(0x1F, 0x4D, 0x78)
AZUL_MEDIO  = RGBColor(0x2E, 0x74, 0xB5)
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)

TNR = 'Times New Roman'

def set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)

def apa_page_numbers(doc):
    """Número de página arriba a la derecha en todas las secciones."""
    for section in doc.sections:
        header = section.header
        p = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
        p.clear()
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run = p.add_run()
        run.font.name = TNR
        run.font.size = Pt(12)
        fldChar1 = OxmlElement('w:fldChar')
        fldChar1.set(qn('w:fldCharType'), 'begin')
        instrText = OxmlElement('w:instrText')
        instrText.set(qn('xml:space'), 'preserve')
        instrText.text = ' PAGE '
        fldChar2 = OxmlElement('w:fldChar')
        fldChar2.set(qn('w:fldCharType'), 'end')
        run._r.append(fldChar1)
        run._r.append(instrText)
        run._r.append(fldChar2)

def apa_fmt(p, indent=True):
    """Aplica doble espacio + sangría APA al párrafo."""
    fmt = p.paragraph_format
    fmt.line_spacing_rule = WD_LINE_SPACING.DOUBLE
    fmt.space_before = Pt(0)
    fmt.space_after  = Pt(0)
    if indent:
        fmt.first_line_indent = Cm(1.27)

# APA Level 1: centrado, negrita
def add_h1(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt(12)
    run.font.name = TNR
    apa_fmt(p, indent=False)
    return p

# APA Level 2: izquierda, negrita
def add_h2(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt(12)
    run.font.name = TNR
    apa_fmt(p, indent=False)
    return p

# APA Level 3: izquierda, negrita cursiva
def add_h3(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(text)
    run.bold = True
    run.italic = True
    run.font.size = Pt(12)
    run.font.name = TNR
    apa_fmt(p, indent=False)
    return p

def add_body(doc, text):
    for line in text.split('\n'):
        p = doc.add_paragraph(line)
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        for run in p.runs:
            run.font.size = Pt(12)
            run.font.name = TNR
        apa_fmt(p, indent=True)
    return p

def add_table(doc, data, has_header=True):
    """Tabla APA: sin color en encabezado, solo líneas horizontales."""
    table = doc.add_table(rows=len(data), cols=len(data[0]))
    table.style = 'Table Grid'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, row_data in enumerate(data):
        row = table.rows[i]
        for j, cell_text in enumerate(row_data):
            cell = row.cells[j]
            cell.text = cell_text
            for para in cell.paragraphs:
                para.alignment = WD_ALIGN_PARAGRAPH.LEFT
                for run in para.runs:
                    run.font.size = Pt(10)
                    run.font.name = TNR
                    if i == 0 and has_header:
                        run.bold = True
    # Quitar bordes verticales internos (estilo APA: solo líneas horizontales)
    for row in table.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            tcBorders = OxmlElement('w:tcBorders')
            for side in ('left', 'right'):
                border = OxmlElement(f'w:{side}')
                border.set(qn('w:val'), 'none')
                tcBorders.append(border)
            tcPr.append(tcBorders)
    # Nota bajo la tabla
    p_note = doc.add_paragraph()
    apa_fmt(p_note, indent=False)
    return table

doc = Document()

# ── Márgenes APA: 2.54 cm todos los lados ────────────────────────
for section in doc.sections:
    section.top_margin    = Cm(2.54)
    section.bottom_margin = Cm(2.54)
    section.left_margin   = Cm(2.54)
    section.right_margin  = Cm(2.54)

apa_page_numbers(doc)

# ── PORTADA APA (student paper) ───────────────────────────────────
# ~3-4 líneas en blanco antes del título
for _ in range(4):
    p = doc.add_paragraph()
    apa_fmt(p, indent=False)

# Logo TECSUP (pequeño, centrado)
if os.path.exists(LOGO):
    p_logo = doc.add_paragraph()
    p_logo.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_logo.paragraph_format.space_before = Pt(0)
    p_logo.paragraph_format.space_after  = Pt(0)
    run_logo = p_logo.add_run()
    run_logo.add_picture(LOGO, width=Cm(5))

p_sp = doc.add_paragraph()
apa_fmt(p_sp, indent=False)

# Título: negrita, centrado
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = p.add_run("Plan de Marketing: LegendaryClass")
r.bold = True; r.font.size = Pt(12); r.font.name = TNR
apa_fmt(p, indent=False)

p2 = doc.add_paragraph()
p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
r2 = p2.add_run("Capítulos 1 al 4")
r2.font.size = Pt(12); r2.font.name = TNR
apa_fmt(p2, indent=False)

p_sp2 = doc.add_paragraph()
apa_fmt(p_sp2, indent=False)

# Autores
for nombre in INTEGRANTES:
    p = doc.add_paragraph(nombre)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in p.runs:
        run.font.size = Pt(12); run.font.name = TNR
    apa_fmt(p, indent=False)

p_sp3 = doc.add_paragraph()
apa_fmt(p_sp3, indent=False)

# Afiliación / Institución
for linea in [CARRERA, "TECSUP"]:
    p = doc.add_paragraph(linea)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in p.runs:
        run.font.size = Pt(12); run.font.name = TNR
    apa_fmt(p, indent=False)

p_sp4 = doc.add_paragraph()
apa_fmt(p_sp4, indent=False)

# Curso, docente, fecha
for linea in [CURSO, DOCENTE, ANIO]:
    p = doc.add_paragraph(linea)
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in p.runs:
        run.font.size = Pt(12); run.font.name = TNR
    apa_fmt(p, indent=False)

doc.add_page_break()

# ── CAP 1 ─────────────────────────────────────────────────────────
add_h1(doc, "Capítulo 1: Introducción")
for titulo, texto in CAP1["secciones"]:
    add_h2(doc, titulo)
    add_body(doc, texto)

doc.add_page_break()

# ── CAP 2 ─────────────────────────────────────────────────────────
add_h1(doc, "Capítulo 2: Análisis del Macroentorno")
add_h2(doc, "Análisis PESTEL")
for sub, texto in CAP2["subsecciones"]:
    add_h3(doc, sub)
    add_body(doc, texto)
add_h2(doc, "Impacto del Macroentorno en el Producto")
add_body(doc, CAP2["impacto"])

doc.add_page_break()

# ── CAP 3 ─────────────────────────────────────────────────────────
add_h1(doc, "Capítulo 3: Segmentación y Perfil del Consumidor")
add_h2(doc, "Segmentación Geográfica")
add_table(doc, [["Nivel","Zona","Justificación"]] + CAP3["geo"])
add_h2(doc, "Segmentación Demográfica")
add_body(doc, CAP3["demo"])
add_h2(doc, "Segmentación Psicográfica")
add_body(doc, CAP3["psico"])
add_h2(doc, "Segmentación Conductual")
add_table(doc, [["Variable","Descripción"]] + CAP3["conductual"])
add_h2(doc, "Perfil del Consumidor Objetivo (Buyer Persona)")
add_body(doc, CAP3["persona"])

doc.add_page_break()

# ── CAP 4 ─────────────────────────────────────────────────────────
add_h1(doc, "Capítulo 4: Estrategias de Marketing Mix")

add_h2(doc, "Producto")
add_h3(doc, "Logo y Slogan")
add_body(doc, 'Logo: Escudo medieval con espada y libro entrelazados, colores morado, dorado y negro.\nSlogan: "Donde el aprendizaje se convierte en leyenda."')
add_h3(doc, "Características y Diferenciadores")
add_table(doc, [["Característica","Beneficio para el colegio"]] + CAP4["producto_feats"])
add_h3(doc, "Tabla Comparativa vs. Competencia")
add_table(doc, CAP4["comparativa"])
add_h3(doc, "Propuesta de Valor")
p_pv = doc.add_paragraph(CAP4["propuesta"])
p_pv.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
p_pv.paragraph_format.left_indent  = Cm(2.54)
p_pv.paragraph_format.right_indent = Cm(2.54)
apa_fmt(p_pv, indent=False)
for run in p_pv.runs:
    run.font.size = Pt(12); run.font.italic = True; run.font.name = TNR

add_h2(doc, "Plaza")
add_body(doc, "Canal directo B2B: visitas y videollamadas con directores académicos.\n"
              "Canal digital: landing page con demo interactiva y formulario de prueba gratuita de 30 días.\n"
              "Canal indirecto (fase 2): alianzas con distribuidoras de tecnología educativa.\n"
              "Cobertura: Fase 1 (meses 1–6) Lima Metropolitana; Fase 2 (meses 7–12) Arequipa y Trujillo.\n"
              "Logística 100% digital: implementación remota con capacitación en dos sesiones virtuales o una presencial.")

add_h2(doc, "Promoción")
add_body(doc, "BTL: ferias educativas, demos presenciales en colegios con piloto gratuito de 30 días, talleres para docentes.\n"
              "Digital: LinkedIn (directores/coordinadores), Facebook/Instagram (docentes y padres), YouTube (demostraciones),\n"
              "email marketing segmentado y WhatsApp Business para seguimiento comercial.")

add_h2(doc, "Precio")
add_h3(doc, "Modelo 1 – Suscripción SaaS")
add_table(doc, CAP4["precio_saas"])
add_h3(doc, "Modelo 2 – Implementación Personalizada (Venta Directa)")
add_table(doc, CAP4["precio_vd"])
add_h3(doc, "Tabla de Costos Estimados")
add_table(doc, CAP4["costos"])
add_h3(doc, "Punto de Equilibrio")
add_body(doc, "Con costos operativos de S/. 1,355/mes, el punto de equilibrio se alcanza con:\n"
              "ocho colegios en el Plan Básico (8 × S/. 180 = S/. 1,440/mes), o\n"
              "cuatro colegios en el Plan Estándar (4 × S/. 350 = S/. 1,400/mes).")
add_h3(doc, "Justificación del Precio")
add_body(doc, "Los precios están posicionados por debajo de Classcraft (aproximadamente S/. 15–20 por alumno al mes). "
              "La estrategia es de valor percibido: el precio se justifica por el retorno en motivación "
              "estudiantil, el ahorro de tiempo docente y la diferenciación institucional. "
              "El descuento anual del 10% incentiva la retención y mejora el flujo de caja del proveedor.")

# ── REFERENCIAS ───────────────────────────────────────────────────
doc.add_page_break()
add_h1(doc, "Referencias")
refs = [
    "HolonIQ. (2023). Global EdTech Market Size 2023. HolonIQ Intelligence.",
    "Instituto Nacional de Estadística e Informática. (2023). Estadísticas de tecnologías de información y comunicación en los hogares. INEI.",
    "Kotler, P., & Keller, K. L. (2022). Marketing management (16th ed.). Pearson Education.",
    "Ministerio de Educación del Perú. (2023). Estadística de la calidad educativa. MINEDU. https://escale.minedu.gob.pe",
    "Osterwalder, A., & Pigneur, Y. (2010). Business model generation. John Wiley & Sons.",
]
for ref in refs:
    p = doc.add_paragraph()
    p.paragraph_format.left_indent        = Cm(1.27)
    p.paragraph_format.first_line_indent  = Cm(-1.27)   # sangría francesa
    p.paragraph_format.line_spacing_rule  = WD_LINE_SPACING.DOUBLE
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after  = Pt(0)
    r = p.add_run(ref)
    r.font.name = TNR; r.font.size = Pt(12)

doc.save(os.path.join(BASE, "LegendaryClass_PlanMarketing.docx"))
print("DOCX OK")

# ═════════════════════════════════════════════════════════════════
# 2. PDF  –  Formato APA 7ª edición
# ═════════════════════════════════════════════════════════════════
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm, inch
from reportlab.lib.colors import HexColor, white, black
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                 TableStyle, PageBreak, Image as RLImage,
                                 Frame, BaseDocTemplate, PageTemplate)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY, TA_RIGHT
from reportlab.platypus.flowables import HRFlowable

PAGE_W, PAGE_H = A4
ML = MR = MB = 2.54*cm
MT = 2.54*cm

def ps(name, **kw):
    base = kw.pop('base', 'Normal')
    return ParagraphStyle(name, fontName=kw.pop('fontName','Times-Roman'),
                          leading=kw.pop('leading', 24),   # doble espacio (12*2)
                          fontSize=kw.pop('fontSize', 12),
                          **kw)

sty_body  = ps("body",  alignment=TA_JUSTIFY, firstLineIndent=1.27*cm,
               spaceAfter=0, spaceBefore=0)
sty_body_no_indent = ps("body_ni", alignment=TA_JUSTIFY,
               spaceAfter=0, spaceBefore=0)
# APA H1: centrado, negrita
sty_h1    = ps("h1",   fontName="Times-Bold",  alignment=TA_CENTER,
               spaceAfter=0, spaceBefore=12)
# APA H2: izquierda, negrita
sty_h2    = ps("h2",   fontName="Times-Bold",  alignment=TA_LEFT,
               spaceAfter=0, spaceBefore=12)
# APA H3: izquierda, negrita-cursiva
sty_h3    = ps("h3",   fontName="Times-BoldItalic", alignment=TA_LEFT,
               spaceAfter=0, spaceBefore=0)

sty_ctr   = ps("ctr",  alignment=TA_CENTER, spaceAfter=0, spaceBefore=0)
sty_title = ps("tt",   fontName="Times-Bold", fontSize=12, alignment=TA_CENTER,
               spaceAfter=0, spaceBefore=0)
sty_ref   = ps("ref",  leftIndent=1.27*cm, firstLineIndent=-1.27*cm,
               spaceAfter=0, spaceBefore=0)
sty_italic= ps("it",   fontName="Times-Italic", alignment=TA_JUSTIFY,
               leftIndent=2.54*cm, rightIndent=2.54*cm,
               spaceAfter=0, spaceBefore=0)
sty_tbl   = ps("tbl",  fontSize=10, leading=14, alignment=TA_LEFT, spaceAfter=0)
sty_tbl_h = ps("tblh", fontName="Times-Bold", fontSize=10, leading=14,
               alignment=TA_LEFT, spaceAfter=0)

LINE = 24  # puntos de interlineado doble para 12pt

def sp():
    return Spacer(1, LINE)   # un salto de línea equivalente a doble espacio

def pdf_table(data, col_widths=None, last_bold=False):
    """Tabla estilo APA: encabezado en negrita, solo líneas horizontales."""
    avail = PAGE_W - ML - MR
    if col_widths is None:
        w = avail / len(data[0])
        col_widths = [w] * len(data[0])

    # Convertir textos a Paragraph para control de fuente
    tdata = []
    for i, row in enumerate(data):
        trow = []
        for cell in row:
            sty = sty_tbl_h if i == 0 else sty_tbl
            trow.append(Paragraph(str(cell), sty))
        tdata.append(trow)

    t = Table(tdata, colWidths=col_widths, repeatRows=1)
    style = [
        # solo líneas horizontales
        ("LINEABOVE",  (0,0),(-1,0),  1, black),
        ("LINEBELOW",  (0,0),(-1,0),  1, black),
        ("LINEBELOW",  (0,-1),(-1,-1),1, black),
        ("VALIGN",     (0,0),(-1,-1), "TOP"),
        ("BOTTOMPADDING",(0,0),(-1,-1), 3),
        ("TOPPADDING",   (0,0),(-1,-1), 3),
    ]
    if last_bold:
        style.append(("FONTNAME",(0,-1),(-1,-1),"Times-Bold"))
    t.setStyle(TableStyle(style))
    return t

# ── Numeración de página (callback) ──────────────────────────────
def _page_num(canvas, doc):
    canvas.saveState()
    canvas.setFont("Times-Roman", 12)
    canvas.drawRightString(PAGE_W - MR, PAGE_H - 1.5*cm, str(doc.page))
    canvas.restoreState()

doc_pdf = BaseDocTemplate(
    os.path.join(BASE, "LegendaryClass_PlanMarketing.pdf"),
    pagesize=A4,
    rightMargin=MR, leftMargin=ML,
    topMargin=MT + 0.8*cm,   # espacio extra para el número de página
    bottomMargin=MB,
)
frame = Frame(ML, MB, PAGE_W - ML - MR, PAGE_H - MT - MB - 0.8*cm, id='main')
template = PageTemplate(id='apa', frames=frame, onPage=_page_num)
doc_pdf.addPageTemplates([template])

story = []

# ── PORTADA APA ───────────────────────────────────────────────────
for _ in range(3):
    story.append(sp())

if os.path.exists(LOGO):
    img = RLImage(LOGO, width=5*cm, height=1.7*cm)
    img.hAlign = 'CENTER'
    story.append(img)
    story.append(sp())

story.append(Paragraph("<b>Plan de Marketing: LegendaryClass</b>", sty_title))
story.append(Paragraph("Capítulos 1 al 4", sty_ctr))
story.append(sp())
for n in INTEGRANTES:
    story.append(Paragraph(n, sty_ctr))
story.append(sp())
for linea in [CARRERA, "TECSUP"]:
    story.append(Paragraph(linea, sty_ctr))
story.append(sp())
for linea in [CURSO, DOCENTE, ANIO]:
    story.append(Paragraph(linea, sty_ctr))
story.append(PageBreak())

# ── CAP 1 ─────────────────────────────────────────────────────────
story.append(Paragraph("<b>Capítulo 1: Introducción</b>", sty_h1))
story.append(sp())
for titulo, texto in CAP1["secciones"]:
    story.append(Paragraph(f"<b>{titulo}</b>", sty_h2))
    story.append(sp())
    for linea in texto.split('\n'):
        story.append(Paragraph(linea, sty_body))
        story.append(sp())
story.append(PageBreak())

# ── CAP 2 ─────────────────────────────────────────────────────────
story.append(Paragraph("<b>Capítulo 2: Análisis del Macroentorno</b>", sty_h1))
story.append(sp())
story.append(Paragraph("<b>Análisis PESTEL</b>", sty_h2))
story.append(sp())
for sub, texto in CAP2["subsecciones"]:
    story.append(Paragraph(f"<i><b>{sub}</b></i>", sty_h3))
    story.append(sp())
    for linea in texto.split('\n'):
        story.append(Paragraph(linea, sty_body))
        story.append(sp())
story.append(Paragraph("<b>Impacto del Macroentorno en el Producto</b>", sty_h2))
story.append(sp())
story.append(Paragraph(CAP2["impacto"], sty_body))
story.append(PageBreak())

# ── CAP 3 ─────────────────────────────────────────────────────────
story.append(Paragraph("<b>Capítulo 3: Segmentación y Perfil del Consumidor</b>", sty_h1))
story.append(sp())
story.append(Paragraph("<b>Segmentación Geográfica</b>", sty_h2))
story.append(sp())
story.append(pdf_table([["Nivel","Zona","Justificación"]] + CAP3["geo"],
                       [3*cm, 4.5*cm, PAGE_W-ML-MR-7.5*cm]))
story.append(sp())
story.append(Paragraph("<b>Segmentación Demográfica</b>", sty_h2))
story.append(sp())
for linea in CAP3["demo"].split('\n'):
    story.append(Paragraph(linea, sty_body))
    story.append(sp())
story.append(Paragraph("<b>Segmentación Psicográfica</b>", sty_h2))
story.append(sp())
for linea in CAP3["psico"].split('\n'):
    story.append(Paragraph(linea, sty_body))
    story.append(sp())
story.append(Paragraph("<b>Segmentación Conductual</b>", sty_h2))
story.append(sp())
story.append(pdf_table([["Variable","Descripción"]] + CAP3["conductual"],
                       [4.5*cm, PAGE_W-ML-MR-4.5*cm]))
story.append(sp())
story.append(Paragraph("<b>Perfil del Consumidor Objetivo (Buyer Persona)</b>", sty_h2))
story.append(sp())
story.append(Paragraph(CAP3["persona"], sty_body))
story.append(PageBreak())

# ── CAP 4 ─────────────────────────────────────────────────────────
story.append(Paragraph("<b>Capítulo 4: Estrategias de Marketing Mix</b>", sty_h1))
story.append(sp())

story.append(Paragraph("<b>Producto</b>", sty_h2)); story.append(sp())
story.append(Paragraph("<i><b>Logo y Slogan</b></i>", sty_h3)); story.append(sp())
story.append(Paragraph('Logo: Escudo medieval con espada y libro entrelazados, colores morado, dorado y negro. Slogan: "Donde el aprendizaje se convierte en leyenda."', sty_body)); story.append(sp())
story.append(Paragraph("<i><b>Características y Diferenciadores</b></i>", sty_h3)); story.append(sp())
story.append(pdf_table([["Característica","Beneficio para el colegio"]] + CAP4["producto_feats"],
                       [8*cm, PAGE_W-ML-MR-8*cm]))
story.append(sp())
story.append(Paragraph("<i><b>Tabla Comparativa vs. Competencia</b></i>", sty_h3)); story.append(sp())
story.append(pdf_table(CAP4["comparativa"],
                       [4*cm, 3.5*cm, 3.5*cm, PAGE_W-ML-MR-11*cm]))
story.append(sp())
story.append(Paragraph("<i><b>Propuesta de Valor</b></i>", sty_h3)); story.append(sp())
story.append(Paragraph(f'<i>{CAP4["propuesta"]}</i>', sty_italic)); story.append(sp())

story.append(Paragraph("<b>Plaza</b>", sty_h2)); story.append(sp())
story.append(Paragraph(
    "Canal directo B2B: visitas y videollamadas con directores académicos. "
    "Canal digital: landing page con demo interactiva y formulario de prueba gratuita de 30 días. "
    "Canal indirecto (fase 2): alianzas con distribuidoras de tecnología educativa. "
    "Cobertura: Fase 1 (meses 1–6) Lima Metropolitana; Fase 2 (meses 7–12) Arequipa y Trujillo. "
    "Logística 100% digital: implementación remota con capacitación en dos sesiones virtuales o una presencial.",
    sty_body)); story.append(sp())

story.append(Paragraph("<b>Promoción</b>", sty_h2)); story.append(sp())
story.append(Paragraph(
    "BTL: ferias educativas, demos presenciales en colegios con piloto gratuito de 30 días, talleres para docentes. "
    "Digital: LinkedIn (directores/coordinadores), Facebook/Instagram (docentes y padres), YouTube (demostraciones), "
    "email marketing segmentado y WhatsApp Business para seguimiento comercial.",
    sty_body)); story.append(sp())

story.append(Paragraph("<b>Precio</b>", sty_h2)); story.append(sp())
story.append(Paragraph("<i><b>Modelo 1 – Suscripción SaaS</b></i>", sty_h3)); story.append(sp())
story.append(pdf_table(CAP4["precio_saas"],
                       [3.5*cm, 4*cm, 3.5*cm, PAGE_W-ML-MR-11*cm]))
story.append(sp())
story.append(Paragraph("<i><b>Modelo 2 – Implementación Personalizada</b></i>", sty_h3)); story.append(sp())
story.append(pdf_table(CAP4["precio_vd"],
                       [5*cm, 3*cm, PAGE_W-ML-MR-8*cm]))
story.append(sp())
story.append(Paragraph("<i><b>Tabla de Costos Estimados</b></i>", sty_h3)); story.append(sp())
story.append(pdf_table(CAP4["costos"],
                       [8.5*cm, 3*cm, PAGE_W-ML-MR-11.5*cm], last_bold=True))
story.append(sp())
story.append(Paragraph("<i><b>Punto de Equilibrio</b></i>", sty_h3)); story.append(sp())
story.append(Paragraph(
    "Con costos operativos de S/. 1,355/mes, el punto de equilibrio se alcanza con "
    "ocho colegios en el Plan Básico (8 × S/. 180 = S/. 1,440/mes) o "
    "cuatro colegios en el Plan Estándar (4 × S/. 350 = S/. 1,400/mes).", sty_body))
story.append(sp())
story.append(Paragraph("<i><b>Justificación del Precio</b></i>", sty_h3)); story.append(sp())
story.append(Paragraph(
    "Los precios están posicionados por debajo de Classcraft (aproximadamente S/. 15–20 por alumno al mes). "
    "La estrategia es de valor percibido: el precio se justifica por el retorno en motivación "
    "estudiantil, el ahorro de tiempo docente y la diferenciación institucional. "
    "El descuento anual del 10% incentiva la retención y mejora el flujo de caja.",
    sty_body))

# ── REFERENCIAS ───────────────────────────────────────────────────
story.append(PageBreak())
story.append(Paragraph("<b>Referencias</b>", sty_h1))
story.append(sp())
pdf_refs = [
    "HolonIQ. (2023). <i>Global EdTech Market Size 2023</i>. HolonIQ Intelligence.",
    "Instituto Nacional de Estadística e Informática. (2023). <i>Estadísticas de tecnologías de información y comunicación en los hogares</i>. INEI.",
    "Kotler, P., &amp; Keller, K. L. (2022). <i>Marketing management</i> (16th ed.). Pearson Education.",
    "Ministerio de Educación del Perú. (2023). <i>Estadística de la calidad educativa</i>. MINEDU. https://escale.minedu.gob.pe",
    "Osterwalder, A., &amp; Pigneur, Y. (2010). <i>Business model generation</i>. John Wiley &amp; Sons.",
]
for ref in pdf_refs:
    story.append(Paragraph(ref, sty_ref))
    story.append(sp())

doc_pdf.build(story)
print("PDF OK")

# ═════════════════════════════════════════════════════════════════
# 3. PPTX
# ═════════════════════════════════════════════════════════════════
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN

PAZ1 = PRGB(0x1F, 0x4D, 0x78)
PAZ2 = PRGB(0x2E, 0x74, 0xB5)
PAZ3 = PRGB(0xD6, 0xE4, 0xF0)
PAZ4 = PRGB(0xEB, 0xF3, 0xFB)
PBLA = PRGB(0xFF, 0xFF, 0xFF)
PNEG = PRGB(0x0F, 0x0F, 0x0F)
PGRI = PRGB(0xF5, 0xF8, 0xFD)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def set_bg_ppt(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, text, l, t, w, h, size=14, bold=False, color=PBLA,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    return sh

def ppt_slide_header(slide, title, subtitle=None):
    """Franja azul oscuro arriba con título blanco"""
    rect(slide, 0, 0, 13.33, 1.05, PAZ1)
    tb(slide, title, 0.3, 0.1, 12.5, 0.85, size=24, bold=True, color=PBLA)
    if subtitle:
        tb(slide, subtitle, 0.3, 0.82, 12.5, 0.3, size=11, color=PAZ3)

def ppt_footer(slide):
    rect(slide, 0, 7.1, 13.33, 0.4, PAZ2)
    tb(slide, "LegendaryClass  |  Plan de Marketing  |  " + ANIO,
       0.3, 7.13, 12.7, 0.3, size=9, color=PBLA, align=PP_ALIGN.CENTER)

def ppt_bullet_box(slide, items, l, t, w, h, title=None, title_color=PAZ1):
    rect(slide, l, t, w, h, PAZ4, PAZ3)
    y_off = 0
    if title:
        tb(slide, title, l+0.12, t+0.08, w-0.2, 0.35, size=11, bold=True, color=title_color)
        y_off = 0.38
    for i, item in enumerate(items):
        tb(slide, item, l+0.12, t+y_off+0.05+i*0.38, w-0.2, 0.35, size=10, color=PNEG)

def ppt_table(slide, data, l, t, w, h):
    rows, cols = len(data), len(data[0])
    col_w = w / cols
    row_h = h / rows
    for i, row in enumerate(data):
        for j, cell in enumerate(row):
            bg = PAZ1 if i == 0 else (PAZ4 if i % 2 == 0 else PBLA)
            fc = PBLA if i == 0 else PNEG
            rect(slide, l + j*col_w, t + i*row_h, col_w, row_h, bg)
            tb(slide, str(cell), l+j*col_w+0.05, t+i*row_h+0.04,
               col_w-0.1, row_h-0.05, size=9,
               bold=(i==0), color=fc, align=PP_ALIGN.CENTER)

# ── SLIDE 1: PORTADA ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg_ppt(s, PBLA)
rect(s, 0, 0, 13.33, 2.0, PAZ1)
rect(s, 0, 6.9, 13.33, 0.6, PAZ2)

if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(0.4), Inches(0.25), width=Inches(3.2))

tb(s, CARRERA, 4.0, 0.2, 9.0, 0.5, size=13, color=PAZ3, align=PP_ALIGN.RIGHT)
tb(s, "PLAN DE MARKETING", 1.0, 2.3, 11.33, 1.0, size=36, bold=True, color=PAZ1, align=PP_ALIGN.CENTER)
tb(s, '"LegendaryClass"', 1.0, 3.3, 11.33, 0.7, size=22, bold=True, color=PAZ2, align=PP_ALIGN.CENTER)
tb(s, "Capítulos 1 al 4", 1.0, 3.95, 11.33, 0.4, size=13, color=PAZ2, align=PP_ALIGN.CENTER, italic=True)

tb(s, f"Curso: {CURSO}  |  Docente: {DOCENTE}", 1.0, 4.55, 11.33, 0.4, size=11, color=PNEG, align=PP_ALIGN.CENTER)
names = "  |  ".join(INTEGRANTES)
tb(s, names, 0.5, 5.05, 12.3, 0.4, size=10, color=PAZ1, align=PP_ALIGN.CENTER)
tb(s, ANIO, 0.5, 6.95, 12.3, 0.35, size=11, bold=True, color=PBLA, align=PP_ALIGN.CENTER)

# ── SLIDE 2: CAP 1 EMPRESA / MISION / VISION ─────────────────────
s = prs.slides.add_slide(blank)
set_bg_ppt(s, PGRI)
ppt_slide_header(s, "CAPÍTULO 1: INTRODUCCIÓN", "Empresa, Misión, Visión y Valores")
ppt_footer(s)

cols = [
    ("MISIÓN", "Transformar la experiencia educativa mediante gamificación que motive a estudiantes, empodere a docentes y conecte a las familias con el progreso académico de sus hijos."),
    ("VISIÓN",  "Ser la plataforma de gamificación educativa líder en Latinoamérica para 2030, en más de 500 instituciones, reconocida por elevar el rendimiento académico."),
    ("VALORES", "• Innovación\n• Compromiso\n• Accesibilidad\n• Transparencia\n• Colaboración"),
]
for i, (tit, txt) in enumerate(cols):
    x = 0.3 + i * 4.3
    rect(s, x, 1.25, 4.0, 5.4, PAZ1)
    tb(s, tit, x+0.12, 1.35, 3.76, 0.45, size=13, bold=True, color=PAZ3)
    tb(s, txt, x+0.12, 1.85, 3.76, 4.6, size=11, color=PBLA)

# empresa label
tb(s, "Empresa: LegendaryClass", 0.3, 6.8, 8.0, 0.3, size=10, bold=True, color=PAZ1)

# ── SLIDE 3: CAP 1 OBJETIVOS ──────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg_ppt(s, PGRI)
ppt_slide_header(s, "Objetivos del Plan de Marketing")
ppt_footer(s)

rect(s, 0.3, 1.2, 12.7, 1.1, PAZ2)
tb(s, "OBJETIVO GENERAL", 0.5, 1.25, 12.3, 0.38, size=12, bold=True, color=PBLA)
tb(s, "Diseñar una estrategia de marketing integral para el lanzamiento de LegendaryClass como plataforma B2B de gamificación educativa en el Perú.", 0.5, 1.6, 12.3, 0.6, size=11, color=PNEG)

esps = [
    "1.  Identificar y segmentar el mercado de colegios privados en Lima Metropolitana con propuestas de valor diferenciadas.",
    "2.  Establecer una estrategia de precios dual (SaaS + implementación personalizada) competitiva frente a alternativas extranjeras.",
    "3.  Diseñar un plan de comunicación digital para directores y coordinadores que genere leads calificados en el primer año.",
]
for i, txt in enumerate(esps):
    rect(s, 0.3, 2.5 + i*1.4, 12.7, 1.25, PAZ4, PAZ3)
    tb(s, txt, 0.5, 2.55 + i*1.4, 12.3, 1.1, size=11, color=PNEG)

# ── SLIDE 4: PESTEL ───────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg_ppt(s, PGRI)
ppt_slide_header(s, "CAPÍTULO 2: Análisis PESTEL")
ppt_footer(s)

celdas = [
    ("POLÍTICO",     "• MINEDU impulsa transformación digital\n• Programas de dotación tecnológica (PRONIED)\n• Ley 29733: protección de datos de menores"),
    ("ECONÓMICO",    "• EdTech global: USD 142B (2023)\n• Post-pandemia: mayor inversión en tecnología\n• Tipo de cambio favorece solución local en soles"),
    ("SOCIAL",       "• Gen Z y Alpha: nativos digitales / gamers\n• Creciente desmotivación estudiantil\n• Padres exigen mayor visibilidad del progreso"),
    ("TECNOLÓGICO",  "• +80% penetración de internet en Lima\n• ClassDojo y Classcraft validan la demanda\n• Stack moderno: Laravel 11, MongoDB"),
    ("ECOLÓGICO",    "• Solución 100% digital: elimina papel\n• Infraestructura cloud con opción verde\n• Posicionamiento eco-friendly"),
    ("LEGAL",        "• Cumplimiento Ley 29733 (menores)\n• Registro IP en INDECOPI\n• Contratos SaaS con SLA y privacidad"),
]
pos = [(0.25,1.15),(4.55,1.15),(8.85,1.15),(0.25,4.05),(4.55,4.05),(8.85,4.05)]
for i, (tit, txt) in enumerate(celdas):
    x, y = pos[i]
    rect(s, x, y, 4.15, 2.65, PAZ1)
    tb(s, tit, x+0.1, y+0.07, 3.95, 0.4, size=12, bold=True, color=PAZ3)
    tb(s, txt, x+0.1, y+0.5,  3.95, 2.05, size=10, color=PBLA)

# ── SLIDE 5: IMPACTO MACROENTORNO ────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg_ppt(s, PGRI)
ppt_slide_header(s, "Impacto del Macroentorno")
ppt_footer(s)

ppt_bullet_box(s,
    ["Digitalización educativa acelerada post-pandemia abre mercado inmediato.",
     "Competidores extranjeros cobran en USD: ventaja para solución local en soles.",
     "Gen Z nativa digital: alta receptividad sin curva de adopción.",
     "Ningún competidor local con sistema RPG completo en español."],
    0.25, 1.2, 6.25, 5.6, title="OPORTUNIDADES", title_color=PAZ1)

ppt_bullet_box(s,
    ["Ley 29733: gestión cuidadosa de datos de menores.",
     "Brecha digital limita expansión fuera de Lima en el corto plazo.",
     "Inestabilidad económica puede retrasar inversiones en colegios."],
    6.85, 1.2, 6.25, 5.6, title="RIESGOS A GESTIONAR", title_color=PAZ2)

# ── SLIDE 6: SEGMENTACIÓN ────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg_ppt(s, PGRI)
ppt_slide_header(s, "CAPÍTULO 3: Segmentación de Mercado")
ppt_footer(s)

segs = [
    ("GEOGRÁFICA",    "• Primario: Lima Metropolitana\n• Secundario: Arequipa, Trujillo, Piura\n• Foco: Surco, Miraflores, Los Olivos, SJL"),
    ("DEMOGRÁFICA",   "• Colegios privados: 150–2,000 alumnos\n• Directores: 35–60 años, NSE B/C\n• Estudiantes: 8–17 años"),
    ("PSICOGRÁFICA",  "• Instituciones que buscan diferenciarse\n• Directivos con mentalidad innovadora\n• Maestros frustrados con sistemas manuales"),
    ("CONDUCTUAL",    "• Ya usan ClassDojo o Google Classroom\n• Uso diario y semanal de la plataforma\n• Alta lealtad post-adopción"),
]
pos2 = [(0.25,1.15),(6.85,1.15),(0.25,4.1),(6.85,4.1)]
for i, (tit, txt) in enumerate(segs):
    x, y = pos2[i]
    rect(s, x, y, 6.15, 2.75, PAZ1)
    tb(s, tit, x+0.12, y+0.1, 5.9, 0.38, size=12, bold=True, color=PAZ3)
    tb(s, txt, x+0.12, y+0.52, 5.9, 2.1, size=11, color=PBLA)

# ── SLIDE 7: BUYER PERSONA ────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg_ppt(s, PGRI)
ppt_slide_header(s, "Perfil del Consumidor Objetivo – Buyer Persona")
ppt_footer(s)

rect(s, 0.25, 1.2, 3.5, 5.6, PAZ1)
tb(s, "RODRIGO", 0.25, 2.5, 3.5, 0.6, size=20, bold=True, color=PAZ3, align=PP_ALIGN.CENTER)
tb(s, "45 años\nDirector de Colegio\nPrivado — Surco, Lima\nNSE B/C", 0.25, 3.15, 3.5, 1.5, size=12, color=PBLA, align=PP_ALIGN.CENTER)

info = [
    ("QUÉ BUSCA",       "• Mejorar disciplina sin aumentar carga docente\n• Diferenciar su colegio de la competencia\n• Datos concretos del comportamiento estudiantil"),
    ("QUÉ LE FRUSTRA",  "• Sistemas de puntos en papel que nadie usa\n• Herramientas en inglés que sus maestros no entienden\n• Pagar en dólares por software genérico"),
    ("QUÉ LO CONVENCE", "• Demo en vivo con datos reales\n• Referencia de otro colegio que lo usa\n• Precio en soles con soporte en español\n• Prueba gratuita de 30 días"),
]
for i, (tit, txt) in enumerate(info):
    x = 4.1 + i * 3.05
    rect(s, x, 1.2, 2.9, 5.6, PAZ2)
    tb(s, tit, x+0.1, 1.3, 2.7, 0.4, size=11, bold=True, color=PAZ3)
    tb(s, txt, x+0.1, 1.75, 2.7, 4.8, size=11, color=PBLA)

# ── SLIDE 8: PRODUCTO ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg_ppt(s, PGRI)
ppt_slide_header(s, "CAPÍTULO 4: Producto")
ppt_footer(s)

feats = [
    "Sistema RPG completo: personajes, XP, niveles, misiones, logros",
    "Multi-rol: Director / Maestro / Estudiante / Padre",
    "Tienda de recompensas personalizable por el colegio",
    "Importación masiva de alumnos desde Excel",
    "Reportes y analítica en tiempo real descargables",
    "Whitelabel: personalizable con logo del colegio",
]
for i, txt in enumerate(feats):
    rect(s, 0.25, 1.2 + i*0.88, 7.0, 0.78, PAZ4, PAZ3)
    tb(s, txt, 0.4, 1.25 + i*0.88, 6.8, 0.68, size=11, color=PNEG)

comp = [
    ["Función",               "LegendaryClass", "ClassDojo", "Classcraft"],
    ["Sistema RPG / Niveles", "SÍ – Completo",  "No tiene",  "Básico"],
    ["100% en español",       "SÍ",             "Parcial",   "No – Inglés"],
    ["Precio en soles",       "SÍ",             "Gratuito",  "No – USD"],
    ["Panel de Directores",   "SÍ",             "No",        "No"],
    ["Módulo para Padres",    "SÍ",             "Sí",        "No"],
]
ppt_table(s, comp, 7.55, 1.2, 5.5, 5.5)

# ── SLIDE 9: PROPUESTA DE VALOR ───────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg_ppt(s, PAZ1)
rect(s, 0, 0, 13.33, 1.05, PAZ2)
rect(s, 0, 6.9, 13.33, 0.6, PAZ2)
ppt_footer(s)

tb(s, "PROPUESTA DE VALOR", 0, 0.1, 13.33, 0.85, size=24, bold=True, color=PBLA, align=PP_ALIGN.CENTER)
tb(s,
   '"LegendaryClass convierte el salón de clases en una experiencia épica de aprendizaje, '
   'incrementando el compromiso y la motivación estudiantil, mientras otorga a los docentes '
   'una herramienta ágil para gestionar comportamientos y reconocer el esfuerzo, '
   'todo en español y al precio de la realidad peruana."',
   1.0, 1.5, 11.33, 4.5, size=22, italic=True, color=PBLA, align=PP_ALIGN.CENTER)

# ── SLIDE 10: PRECIO ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg_ppt(s, PGRI)
ppt_slide_header(s, "Estrategia de Precio", "Modelo SaaS + Implementación Personalizada")
ppt_footer(s)

tb(s, "Modelo 1 — Suscripción SaaS", 0.25, 1.15, 6.3, 0.4, size=12, bold=True, color=PAZ1)
ppt_table(s, CAP4["precio_saas"], 0.25, 1.6, 6.3, 2.5)

tb(s, "Modelo 2 — Implementación Personalizada", 7.0, 1.15, 6.1, 0.4, size=12, bold=True, color=PAZ1)
ppt_table(s, CAP4["precio_vd"], 7.0, 1.6, 6.1, 2.0)

tb(s, "Tabla de Costos Estimados", 0.25, 4.3, 12.8, 0.4, size=12, bold=True, color=PAZ1)
ppt_table(s, CAP4["costos"], 0.25, 4.75, 12.8, 2.0)

# ── SLIDE 11: CIERRE ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg_ppt(s, PGRI)
ppt_slide_header(s, "Resumen y Próximos Pasos")
ppt_footer(s)

caps = [
    ("Cap. 1", "Introducción\nEmpresa · Misión\nVisión · Objetivos"),
    ("Cap. 2", "Macroentorno\nAnálisis PESTEL\nImpacto y oportunidades"),
    ("Cap. 3", "Segmentación\nGeográfica / Demo\nBuyer Persona"),
    ("Cap. 4", "Marketing Mix\nProducto · Propuesta\nde Valor · Precio"),
]
for i, (num, texto) in enumerate(caps):
    x = 0.25 + i * 3.2
    rect(s, x, 1.2, 3.05, 4.0, PAZ1)
    tb(s, "✓  " + num, x+0.1, 1.3, 2.85, 0.55, size=16, bold=True, color=PAZ3, align=PP_ALIGN.CENTER)
    tb(s, texto, x+0.1, 1.95, 2.85, 3.0, size=12, color=PBLA, align=PP_ALIGN.CENTER)

rect(s, 0.25, 5.5, 12.83, 1.0, PAZ2)
tb(s, "Próximos avances: Cap. 5 Investigación de Mercado  ·  Cap. 6 Marketing Digital  ·  Estrategia Promoción completa",
   0.4, 5.6, 12.5, 0.8, size=12, bold=True, color=PBLA, align=PP_ALIGN.CENTER)

prs.save(os.path.join(BASE, "LegendaryClass_Lab01.pptx"))
print("PPTX OK")
print("\nListo: DOCX, PDF y PPTX generados con estilo TECSUP.")
