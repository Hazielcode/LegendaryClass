from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

MORADO     = RGBColor(0x4C, 0x1D, 0x95)
DORADO     = RGBColor(0xD9, 0x7B, 0x06)
NEGRO      = RGBColor(0x0F, 0x0F, 0x0F)
BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLARO = RGBColor(0xF3, 0xF0, 0xFF)
MORADO_MED = RGBColor(0x6D, 0x28, 0xD9)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, l, t, w, h, size=18, bold=False, color=BLANCO, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── SLIDE 1: PORTADA ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, MORADO)
add_rect(s, 0, 6.4, 13.33, 1.1, DORADO)
add_textbox(s, "LegendaryClass", 1, 1.0, 11, 1.8, size=60, bold=True, color=DORADO, align=PP_ALIGN.CENTER)
add_textbox(s, '"Donde el aprendizaje se convierte en leyenda"', 1, 2.7, 11, 0.7, size=20, italic=True, color=BLANCO, align=PP_ALIGN.CENTER)
add_rect(s, 3.5, 3.5, 6.33, 0.04, DORADO)
integrantes = [
    "Aguirre Saavedra, Juan Alexis",
    "Alfonso Solorzano, Samir Haziel",
    "Galvan Morales, Luis Enrique",
    "Galvan Guerrero, Matias",
]
for i, nombre in enumerate(integrantes):
    add_textbox(s, nombre, 2, 3.75 + i*0.38, 9.33, 0.45, size=14, color=BLANCO, align=PP_ALIGN.CENTER)
add_textbox(s, "Marketing y Comercializacion de Nuevos Productos  |  Mg. Miluska Horna Elera  |  2026", 1, 6.6, 11.33, 0.5, size=13, color=MORADO, align=PP_ALIGN.CENTER, bold=True)

# ── SLIDE 2: QUE ES LEGENDARYCLASS ───────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, GRIS_CLARO)
add_rect(s, 0, 0, 13.33, 1.1, MORADO)
add_textbox(s, "CAPITULO 1 - Que es LegendaryClass?", 0.3, 0.15, 12.5, 0.8, size=26, bold=True, color=BLANCO)
bullets = [
    ("Plataforma web de gamificacion educativa", "Transforma el aula en un RPG: personajes, XP, niveles, misiones y recompensas."),
    ("Multi-rol completo", "Director, Maestro, Estudiante y Padre con su panel personalizado."),
    ("Sistema de recompensas y comportamiento", "Docentes registran comportamientos; alumnos acumulan puntos y los canjean en la tienda."),
    ("Reportes y analitica en tiempo real", "Directores y maestros visualizan progreso individual y grupal al instante."),
]
for i, (titulo, desc) in enumerate(bullets):
    y = 1.35 + i * 1.35
    add_rect(s, 0.4, y, 12.5, 1.15, MORADO)
    add_textbox(s, titulo, 0.6, y + 0.05, 12, 0.42, size=16, bold=True, color=DORADO)
    add_textbox(s, desc,   0.6, y + 0.48, 12, 0.55, size=13, color=BLANCO)

# ── SLIDE 3: MISION VISION VALORES ───────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, GRIS_CLARO)
add_rect(s, 0, 0, 13.33, 1.1, MORADO)
add_textbox(s, "Mision - Vision - Valores", 0.3, 0.15, 12.5, 0.8, size=26, bold=True, color=BLANCO)
cols = [
    ("MISION", "Transformar la experiencia educativa mediante gamificacion que motive a estudiantes, empodere a docentes y conecte a las familias con el progreso academico de sus hijos."),
    ("VISION",  "Ser la plataforma de gamificacion educativa lider en Latinoamerica para 2030, presente en mas de 500 instituciones y reconocida por elevar el rendimiento academico."),
    ("VALORES", "- Innovacion\n- Compromiso\n- Accesibilidad\n- Transparencia\n- Colaboracion"),
]
for i, (titulo, texto) in enumerate(cols):
    x = 0.3 + i * 4.3
    add_rect(s, x, 1.3, 4.0, 5.7, MORADO)
    add_textbox(s, titulo, x+0.15, 1.45, 3.7, 0.55, size=15, bold=True, color=DORADO)
    add_textbox(s, texto,  x+0.15, 2.1,  3.7, 4.7, size=13, color=BLANCO)

# ── SLIDE 4: OBJETIVOS ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, GRIS_CLARO)
add_rect(s, 0, 0, 13.33, 1.1, MORADO)
add_textbox(s, "Objetivos", 0.3, 0.15, 12.5, 0.8, size=26, bold=True, color=BLANCO)
add_rect(s, 0.4, 1.2, 12.5, 1.35, DORADO)
add_textbox(s, "OBJETIVO GENERAL", 0.6, 1.25, 12, 0.4, size=14, bold=True, color=MORADO)
add_textbox(s, "Disenar una estrategia de marketing integral para el lanzamiento de LegendaryClass como plataforma de gamificacion educativa B2B, estableciendo modelos de comercializacion sostenibles y escalables en el mercado peruano.", 0.6, 1.62, 12, 0.8, size=13, color=NEGRO)
esp = [
    "1.  Identificar y segmentar el mercado de colegios privados en Lima Metropolitana para definir propuestas de valor diferenciadas por tamano e institucion.",
    "2.  Establecer una estrategia de precios dual (SaaS + implementacion personalizada) competitiva frente a alternativas extranjeras y adaptada a la realidad peruana.",
    "3.  Disenar un plan de comunicacion digital orientado a directores y coordinadores academicos para generar leads calificados en el primer ano.",
]
for i, txt in enumerate(esp):
    y = 2.75 + i * 1.45
    add_rect(s, 0.4, y, 12.5, 1.28, MORADO)
    add_textbox(s, txt, 0.6, y+0.18, 12, 1.0, size=13, color=BLANCO)

# ── SLIDE 5: PESTEL ───────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, GRIS_CLARO)
add_rect(s, 0, 0, 13.33, 1.1, MORADO)
add_textbox(s, "CAPITULO 2 - Analisis PESTEL", 0.3, 0.15, 12.5, 0.8, size=26, bold=True, color=BLANCO)
celdas = [
    ("P - POLITICO",     "- MINEDU impulsa transformacion digital\n- Programas de dotacion tecnologica (PRONIED)\n- Ley de proteccion de datos de menores (Ley 29733)"),
    ("E - ECONOMICO",    "- Mercado EdTech global: USD 142B (2023)\n- Post-pandemia: mayor inversion en tecnologia educativa\n- Tipo de cambio favorece solucion local en soles"),
    ("S - SOCIAL",       "- Gen Z y Alpha: nativos digitales / gamers\n- Creciente desmotivacion estudiantil\n- Padres exigen mayor visibilidad del progreso"),
    ("T - TECNOLOGICO",  "- +80% penetracion de internet en Lima\n- ClassDojo y Classcraft validan la demanda global\n- Stack moderno: Laravel 11, MongoDB, TailwindCSS"),
    ("E - ECOLOGICO",    "- Solucion 100% digital: elimina papel y materiales fisicos\n- Infraestructura cloud con opcion de servidores verdes"),
    ("L - LEGAL",        "- Cumplimiento Ley 29733 (datos de menores)\n- Registro de propiedad intelectual en INDECOPI\n- Contratos SaaS con SLA y clausulas de privacidad"),
]
pos = [(0.3,1.2),(4.55,1.2),(8.8,1.2),(0.3,4.1),(4.55,4.1),(8.8,4.1)]
for i, (titulo, texto) in enumerate(celdas):
    x, y = pos[i]
    add_rect(s, x, y, 4.1, 2.7, MORADO)
    add_textbox(s, titulo, x+0.1, y+0.08, 3.9, 0.42, size=14, bold=True, color=DORADO)
    add_textbox(s, texto,  x+0.1, y+0.52, 3.9, 2.1, size=11, color=BLANCO)

# ── SLIDE 6: IMPACTO MACROENTORNO ────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, GRIS_CLARO)
add_rect(s, 0, 0, 13.33, 1.1, MORADO)
add_textbox(s, "Impacto del Macroentorno", 0.3, 0.15, 12.5, 0.8, size=26, bold=True, color=BLANCO)
ops = [
    "Digitalizacion educativa acelerada post-pandemia abre mercado inmediato.",
    "Competidores extranjeros cobran en USD: ventaja para solucion local en soles.",
    "Gen Z nativa digital: alta receptividad a la gamificacion sin curva de adopcion.",
    "Ningun competidor local con sistema RPG completo en espanol.",
]
ries = [
    "Ley 29733: exige gestion cuidadosa de datos de menores de edad.",
    "Brecha digital limita expansion inmediata fuera de Lima.",
    "Inestabilidad economica puede retrasar inversiones en colegios publicos.",
]
add_rect(s, 0.3, 1.2, 6.1, 5.8, MORADO)
add_textbox(s, "OPORTUNIDADES", 0.5, 1.3, 5.7, 0.45, size=15, bold=True, color=DORADO)
for i, txt in enumerate(ops):
    add_textbox(s, "- " + txt, 0.5, 1.9+i*1.22, 5.7, 1.1, size=12, color=BLANCO)
add_rect(s, 6.9, 1.2, 6.1, 5.8, MORADO_MED)
add_textbox(s, "RIESGOS A GESTIONAR", 7.1, 1.3, 5.7, 0.45, size=15, bold=True, color=DORADO)
for i, txt in enumerate(ries):
    add_textbox(s, "- " + txt, 7.1, 1.9+i*1.55, 5.7, 1.4, size=12, color=BLANCO)

# ── SLIDE 7: SEGMENTACION ────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, GRIS_CLARO)
add_rect(s, 0, 0, 13.33, 1.1, MORADO)
add_textbox(s, "CAPITULO 3 - Segmentacion de Mercado", 0.3, 0.15, 12.5, 0.8, size=26, bold=True, color=BLANCO)
segs = [
    ("GEOGRAFICA",    "- Primario: Lima Metropolitana\n- Secundario: Arequipa, Trujillo, Piura\n- Foco inicial: Surco, Miraflores, Los Olivos, SJL"),
    ("DEMOGRAFICA",   "- Colegios privados: 150-2,000 alumnos\n- Directores: 35-60 anos, NSE B y C\n- Estudiantes: 8-17 anos (primaria alta / secundaria)"),
    ("PSICOGRAFICA",  "- Instituciones que buscan diferenciarse\n- Directivos con mentalidad innovadora\n- Maestros frustrados con sistemas manuales"),
    ("CONDUCTUAL",    "- Colegios que ya usan ClassDojo o Google Classroom\n- Uso diario (comportamientos) y semanal (reportes)\n- Alta lealtad post-adopcion por datos acumulados"),
]
pos2 = [(0.3,1.2),(6.9,1.2),(0.3,4.3),(6.9,4.3)]
for i, (titulo, texto) in enumerate(segs):
    x, y = pos2[i]
    add_rect(s, x, y, 6.1, 2.95, MORADO)
    add_textbox(s, titulo, x+0.15, y+0.1, 5.8, 0.45, size=14, bold=True, color=DORADO)
    add_textbox(s, texto,  x+0.15, y+0.6, 5.8, 2.2, size=12, color=BLANCO)

# ── SLIDE 8: BUYER PERSONA ────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, GRIS_CLARO)
add_rect(s, 0, 0, 13.33, 1.1, MORADO)
add_textbox(s, "Perfil del Consumidor Objetivo", 0.3, 0.15, 12.5, 0.8, size=26, bold=True, color=BLANCO)
add_rect(s, 0.3, 1.2, 3.5, 5.8, MORADO)
add_textbox(s, "RODRIGO", 0.3, 2.8, 3.5, 0.6, size=22, bold=True, color=DORADO, align=PP_ALIGN.CENTER)
add_textbox(s, "45 anos\nDirector de Colegio Privado\nSurco, Lima\nNSE B/C", 0.3, 3.45, 3.5, 1.2, size=13, color=BLANCO, align=PP_ALIGN.CENTER)
cols_p = [
    ("QUE BUSCA",       "- Mejorar disciplina sin aumentar carga docente\n- Diferenciar su colegio frente a competidores\n- Datos concretos del comportamiento estudiantil"),
    ("QUE LE FRUSTRA",  "- Sistemas de puntos en papel que nadie usa\n- Herramientas en ingles que sus maestros no entienden\n- Pagar en dolares por software generico"),
    ("QUE LO CONVENCE", "- Demo en vivo con datos reales\n- Referencia de otro colegio que ya lo usa\n- Precio en soles con soporte en espanol\n- Prueba gratuita de 30 dias"),
]
for i, (titulo, texto) in enumerate(cols_p):
    x = 4.1 + i * 3.05
    add_rect(s, x, 1.2, 2.85, 5.8, MORADO_MED)
    add_textbox(s, titulo, x+0.1, 1.3, 2.65, 0.45, size=13, bold=True, color=DORADO)
    add_textbox(s, texto,  x+0.1, 1.85, 2.65, 4.9, size=12, color=BLANCO)

# ── SLIDE 9: PRODUCTO + COMPARATIVA ──────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, GRIS_CLARO)
add_rect(s, 0, 0, 13.33, 1.1, MORADO)
add_textbox(s, "CAPITULO 4 - Producto", 0.3, 0.15, 12.5, 0.8, size=26, bold=True, color=BLANCO)
feats = [
    "Sistema RPG completo: personajes, XP, niveles, misiones, logros",
    "Multi-rol: Director / Maestro / Estudiante / Padre",
    "Tienda de recompensas personalizable por el colegio",
    "Reportes y analitica en tiempo real descargables",
    "Importacion masiva de alumnos desde Excel",
    "Personalizable con logo y marca del colegio (whitelabel)",
]
for i, txt in enumerate(feats):
    y = 1.25 + i * 0.88
    add_rect(s, 0.3, y, 7.0, 0.75, MORADO)
    add_textbox(s, txt, 0.5, y+0.1, 6.8, 0.55, size=12, color=BLANCO)
headers = ["Funcion", "LegendaryClass", "ClassDojo", "Classcraft"]
rows = [
    ["Sistema RPG / Niveles", "SI - Completo", "NO tiene", "Basico"],
    ["100% en espanol",       "SI",            "Parcial",  "NO - Ingles"],
    ["Precio en soles",       "SI",            "Gratuito", "NO - USD"],
    ["Panel de Directores",   "SI",            "NO",       "NO"],
    ["Modulo para Padres",    "SI",            "SI",       "NO"],
    ["Tienda de Recompensas", "SI - Completa", "NO",       "Limitado"],
]
col_x = [7.6, 9.3, 10.85, 12.1]
col_w = [1.65, 1.5, 1.2, 1.2]
add_rect(s, 7.5, 1.2, 5.6, 0.5, DORADO)
for j, h in enumerate(headers):
    add_textbox(s, h, col_x[j], 1.25, col_w[j], 0.4, size=11, bold=True, color=NEGRO, align=PP_ALIGN.CENTER)
for i, row in enumerate(rows):
    bg = MORADO if i % 2 == 0 else MORADO_MED
    add_rect(s, 7.5, 1.75 + i*0.9, 5.6, 0.85, bg)
    for j, cell in enumerate(row):
        add_textbox(s, cell, col_x[j], 1.82 + i*0.9, col_w[j], 0.72, size=10, color=BLANCO, align=PP_ALIGN.CENTER)

# ── SLIDE 10: PROPUESTA DE VALOR ──────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, MORADO)
add_rect(s, 0, 0, 13.33, 1.1, DORADO)
add_textbox(s, "PROPUESTA DE VALOR", 0, 0.15, 13.33, 0.8, size=28, bold=True, color=MORADO, align=PP_ALIGN.CENTER)
add_textbox(s,
    "LegendaryClass convierte el salon de clases en una\nexperiencia epica de aprendizaje, incrementando el\ncompromiso y la motivacion estudiantil, mientras otorga\na los docentes una herramienta agil para gestionar\ncomportamientos y reconocer el esfuerzo,\ntodo en espanol y al precio de la realidad peruana.",
    1.0, 1.5, 11.33, 4.5, size=24, italic=True, color=BLANCO, align=PP_ALIGN.CENTER)
add_rect(s, 2.0, 6.2, 9.33, 0.8, DORADO)
add_textbox(s, "legendaryclass.com  |  Prueba gratuita 30 dias  |  Soporte en espanol", 2.0, 6.3, 9.33, 0.6, size=14, bold=True, color=MORADO, align=PP_ALIGN.CENTER)

# ── SLIDE 11: CIERRE ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, GRIS_CLARO)
add_rect(s, 0, 0, 13.33, 1.1, MORADO)
add_textbox(s, "Resumen y Proximos Pasos", 0.3, 0.15, 12.5, 0.8, size=26, bold=True, color=BLANCO)
caps = [
    ("Cap. 1", "Introduccion\nEmpresa - Mision\nVision - Objetivos"),
    ("Cap. 2", "Macroentorno\nAnalisis PESTEL\nImpacto y oportunidades"),
    ("Cap. 3", "Segmentacion\nGeografica / Demo\nBuyer Persona"),
    ("Cap. 4", "Marketing Mix\nProducto - Propuesta\nde Valor - Marca"),
]
for i, (num, texto) in enumerate(caps):
    x = 0.3 + i * 3.2
    add_rect(s, x, 1.3, 3.0, 3.5, MORADO)
    add_textbox(s, "OK  " + num, x+0.1, 1.4, 2.8, 0.55, size=18, bold=True, color=DORADO, align=PP_ALIGN.CENTER)
    add_textbox(s, texto, x+0.1, 2.0, 2.8, 2.7, size=13, color=BLANCO, align=PP_ALIGN.CENTER)
add_rect(s, 0.3, 5.1, 12.7, 1.1, DORADO)
add_textbox(s, "Proximos: Cap.5 Investigacion de Mercado  -  Cap.6 Marketing Digital  -  Estrategia de Precio completa", 0.5, 5.22, 12.3, 0.8, size=14, bold=True, color=MORADO, align=PP_ALIGN.CENTER)
add_textbox(s, "Gracias!  -  LegendaryClass", 0.3, 6.35, 12.7, 0.8, size=20, bold=True, color=MORADO, align=PP_ALIGN.CENTER)

prs.save("C:/xampp/htdocs/LegendaryClass/LegendaryClass_Lab01.pptx")
print("PPT generada OK")
